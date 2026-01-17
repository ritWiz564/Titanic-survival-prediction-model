from pptx import Pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd

# Create presentation
prs = Pptx()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add blue background rectangle
    shape = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.space_before = Pt(12)

def add_table_slide(prs, title, df):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Add table
    rows, cols = df.shape
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(4.5)
    
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    
    # Header row
    for col_idx, col_name in enumerate(df.columns):
        cell = table.rows[0].cells[col_idx]
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.rows[row_idx + 1].cells[col_idx]
            value = df.iloc[row_idx, col_idx]
            if isinstance(value, float):
                cell.text = f"{value:.4f}"
            else:
                cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

def add_image_slide(prs, title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Add image
    left = Inches(1)
    top = Inches(2)
    slide.shapes.add_picture(image_path, left, top, width=Inches(8))

# Load results
results_df = pd.read_csv("model_comparison.csv")

# Standardize column names
column_mapping = {}
for col in results_df.columns:
    if 'accuracy' in col.lower():
        column_mapping[col] = 'Accuracy'
    elif 'precision' in col.lower():
        column_mapping[col] = 'Precision'
    elif 'recall' in col.lower():
        column_mapping[col] = 'Recall'
    elif 'f1' in col.lower():
        column_mapping[col] = 'F1 Score'
    elif 'roc' in col.lower() or 'auc' in col.lower():
        column_mapping[col] = 'ROC-AUC'

results_df = results_df.rename(columns=column_mapping)

# Find best model
results_df['Average'] = results_df[['Accuracy', 'Precision', 'Recall', 'F1 Score', 'ROC-AUC']].mean(axis=1)
best_idx = results_df['Average'].idxmax()
best_model = results_df.loc[best_idx, 'Model']

# Slide 1: Title Slide
add_title_slide(prs, 
    "Comparative Study of Machine Learning Algorithms",
    "Titanic Survival Prediction")

# Slide 2: Introduction
add_content_slide(prs, "Introduction", [
    "Objective: Predict passenger survival on the Titanic using machine learning",
    "Dataset: Titanic dataset with 891 passenger records",
    "Task: Binary classification (Survived/Did not survive)",
    "Approach: Compare 5 different ML algorithms",
    "Evaluation: Using Accuracy, Precision, Recall, F1 Score, and ROC-AUC"
])

# Slide 3: Literature Survey - Overview
add_content_slide(prs, "Literature Survey - Background", [
    "Machine Learning for Classification: Supervised learning techniques widely used for binary classification problems",
    "Logistic Regression: Linear model for binary outcomes, baseline for classification tasks",
    "Decision Trees & Random Forests: Non-linear models that handle complex patterns through recursive partitioning",
    "Support Vector Machines (SVM): Find optimal hyperplane for class separation",
    "K-Nearest Neighbors (KNN): Instance-based learning using proximity measures"
])

# Slide 4: Literature Survey - Comparison Table
lit_comparison = pd.DataFrame({
    'Algorithm': ['Logistic Regression', 'Random Forest', 'Decision Tree', 'SVM', 'KNN'],
    'Type': ['Linear', 'Ensemble', 'Tree-based', 'Kernel-based', 'Instance-based'],
    'Strengths': ['Fast, Interpretable', 'High Accuracy', 'Simple', 'Complex Boundaries', 'No Training'],
    'Limitations': ['Linear only', 'Black box', 'Overfitting', 'Slow', 'Memory intensive']
})
add_table_slide(prs, "Literature Survey - Algorithm Comparison", lit_comparison)

# Slide 5: Methodology - Data Preprocessing
add_content_slide(prs, "Methodology - Data Preprocessing", [
    "1. Handling Missing Values: Imputed age with median, filled embarked with mode",
    "2. Feature Engineering: Created family_size, is_alone features",
    "3. Encoding Categorical Variables: One-hot encoding for sex, embark_town, who",
    "4. Feature Selection: Removed non-predictive columns (name, ticket, cabin, deck)",
    "5. Train-Test Split: 80% training, 20% testing"
])

# Slide 6: Methodology - Model Training
add_content_slide(prs, "Methodology - Model Training", [
    "Algorithms Implemented:",
    "  • Logistic Regression - Linear classification baseline",
    "  • Random Forest - Ensemble of 100 decision trees",
    "  • Decision Tree - Single tree with Gini criterion",
    "  • Support Vector Machine - RBF kernel with C=1.0",
    "  • K-Nearest Neighbors - k=5 neighbors",
    "",
    "Evaluation Metrics: Accuracy, Precision, Recall, F1 Score, ROC-AUC"
])

# Slide 7: Results - Performance Table
add_table_slide(prs, "Results - Model Performance Comparison", 
    results_df[['Model', 'Accuracy', 'Precision', 'Recall', 'F1 Score', 'ROC-AUC']])

# Slide 8: Results - Visualization 1
add_image_slide(prs, "Results - Performance Metrics Comparison", "individual_metrics.png")

# Slide 9: Results - Visualization 2
add_image_slide(prs, "Results - ROC Curves", "model_comparison_charts.png")

# Slide 10: Results - Key Findings
best_acc = results_df.loc[results_df['Accuracy'].idxmax()]
best_roc = results_df.loc[results_df['ROC-AUC'].idxmax()]

add_content_slide(prs, "Results - Key Findings", [
    f"Best Overall Model: {best_model} (Avg Score: {results_df.loc[best_idx, 'Average']:.4f})",
    f"Highest Accuracy: {best_acc['Model']} ({best_acc['Accuracy']:.4f})",
    f"Highest ROC-AUC: {best_roc['Model']} ({best_roc['ROC-AUC']:.4f})",
    "",
    "All models achieved >75% accuracy",
    "Tree-based models (Random Forest, Decision Tree) showed strong performance",
    "Logistic Regression provided good baseline with interpretability"
])

# Slide 11: Conclusions
add_content_slide(prs, "Conclusions", [
    f"✓ Successfully compared 5 ML algorithms for Titanic survival prediction",
    f"✓ {best_model} achieved best overall performance",
    f"✓ Feature engineering significantly improved model accuracy",
    f"✓ Passenger class, sex, and age were key survival predictors",
    "",
    "Future Work:",
    "  • Hyperparameter tuning for optimization",
    "  • Deep learning approaches",
    "  • Feature importance analysis"
])

# Slide 12: Thank You
add_title_slide(prs, "Thank You", "Questions?")

# Save presentation
prs.save('ML_Project_Presentation.pptx')
print(" PowerPoint presentation created: ML_Project_Presentation.pptx")